#!/usr/bin/env python3
"""
BMW Brand Skill - AI for POs with GitHub Copilot Presentation
Creates a presentation showcasing the agentSkills repo and how a PO uses
GitHub Copilot with custom skills and MCP servers.
Department: DE-841

Usage:
    python examples/create_ai_for_pos_with_copilot.py

Output:
    AI_for_POs_with_Copilot.pptx in the skill root folder
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
EXAMPLES_DIR = os.path.dirname(os.path.abspath(__file__))
SKILL_ROOT = os.path.dirname(EXAMPLES_DIR)
TEMPLATE_PATH = os.path.join(SKILL_ROOT, "Presentation1.pptx")
OUTPUT_DIR = SKILL_ROOT

# Default flow font (from brand-styling-guide.md)
BMW_FONT = "BMWGroupTN Condensed"
BMW_FONT_FALLBACK = "Arial"

# BMW Brand Colors (BMW Group 21 theme)
BMW_COLORS = {
    'ocean_blue': RGBColor(0x03, 0x59, 0x70),
    'night_blue': RGBColor(0x17, 0x3B, 0x68),
    'black': RGBColor(0x00, 0x00, 0x00),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x92, 0xA2, 0xBD),
    'accent_teal': RGBColor(0x54, 0x8D, 0x9E),
    'blue_gray': RGBColor(0x85, 0xAC, 0xB9),
    'pale_blue': RGBColor(0xAB, 0xC4, 0xCF),
    'ice_blue_4': RGBColor(0xC8, 0xD7, 0xE0),
    'ice_blue_5': RGBColor(0xDE, 0xE5, 0xEC),
    'ice_blue_6': RGBColor(0xE8, 0xEB, 0xF1),
    'cyan': RGBColor(0x07, 0x9E, 0xDA),
    'green': RGBColor(0x50, 0x81, 0x30),
    'orange': RGBColor(0xE9, 0x6D, 0x0C),
    'footer_gray': RGBColor(0x7A, 0x7A, 0x7A),
}


# ── Helper Functions ─────────────────────────────────────────────────────

def _update_master_footer(prs, footer_text):
    """Update the master slide's built-in FußzeileAU1 footer text."""
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.name == 'FußzeileAU1' or shape.name == 'Fu\u00dFzeileAU1':
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = ''
                shape.text_frame.paragraphs[0].runs[0].text = footer_text


def _find_layout_by_name(prs, name):
    """Find a slide layout by name across all slide masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise ValueError(f"Layout '{name}' not found in any slide master")


def _delete_existing_slides(prs):
    """Remove all existing slides, keeping masters/layouts intact."""
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId is None:
            for attr_name in sldId.attrib:
                if attr_name.endswith("}id") or attr_name == "r:id":
                    rId = sldId.attrib[attr_name]
                    break
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


# ── Slide Builders ───────────────────────────────────────────────────────

def create_title_slide(prs, title, subtitle, department="DE-841"):
    """Create a title slide using 'Title | Full Picture' layout."""
    layout = _find_layout_by_name(prs, "Title | Full Picture")
    slide = prs.slides.add_slide(layout)

    nsmap_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'

    for shape in slide.shapes:
        sp = shape._element
        ph_el = sp.find(f'.//{nsmap_p}ph')
        if ph_el is None:
            continue
        ph_type = ph_el.get('type')
        ph_idx = ph_el.get('idx')

        if ph_type == 'ctrTitle':
            shape.text_frame.paragraphs[0].text = title.upper()
            for p in shape.text_frame.paragraphs:
                p.font.name = BMW_FONT
                p.font.size = Pt(40)
                p.font.color.rgb = BMW_COLORS['white']
                p.alignment = PP_ALIGN.LEFT

        elif ph_type == 'subTitle' and ph_idx == '1':
            shape.text_frame.paragraphs[0].text = subtitle.upper()
            for p in shape.text_frame.paragraphs:
                p.font.name = BMW_FONT
                p.font.size = Pt(20)
                p.font.color.rgb = BMW_COLORS['white']
                p.alignment = PP_ALIGN.LEFT

        elif ph_type == 'body' and ph_idx == '22':
            shape.text_frame.paragraphs[0].text = department
            for p in shape.text_frame.paragraphs:
                p.font.name = BMW_FONT
                p.font.size = Pt(18)
                p.font.color.rgb = BMW_COLORS['white']
                p.alignment = PP_ALIGN.LEFT

        elif ph_type in ('pic', 'dgm'):
            pass  # Background image and SmartArt logos inherit from layout

    return slide


def create_content_slide(prs, title, bullet_points):
    """Create a content slide with bullet points using Grid | 1 layout."""
    try:
        layout = _find_layout_by_name(prs, "Grid | 1")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Title — 26pt UPPERCASE Ocean Blue per brand guide
    title_box = slide.shapes.add_textbox(
        Emu(488947), Emu(347184), Emu(11224684), Emu(400110)
    )
    title_box.text_frame.word_wrap = True
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(26)
    title_para.font.bold = False
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']
    title_para.alignment = PP_ALIGN.LEFT

    # Body area
    content_box = slide.shapes.add_textbox(
        Emu(488947), Emu(1413933), Emu(11224684), Emu(4894792)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points[:5]):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = f"\u25aa  {point}"
        para.font.name = BMW_FONT
        para.font.size = Pt(18)
        para.font.color.rgb = BMW_COLORS['black']
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(6)

    return slide


def create_content_with_sub_bullets(prs, title, items):
    """Create a content slide with main bullets and sub-bullets.
    items: list of tuples (main_text, [sub_text, ...])
    """
    try:
        layout = _find_layout_by_name(prs, "Grid | 1")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Emu(488947), Emu(347184), Emu(11224684), Emu(400110)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(26)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    content_box = slide.shapes.add_textbox(
        Emu(488947), Emu(1413933), Emu(11224684), Emu(4894792)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    first = True
    for main_text, sub_items in items:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.text = f"\u25aa  {main_text}"
        para.font.name = BMW_FONT
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = BMW_COLORS['black']
        para.space_before = Pt(8)
        para.space_after = Pt(2)

        for sub in sub_items[:3]:
            sub_para = tf.add_paragraph()
            sub_para.text = f"     -  {sub}"
            sub_para.font.name = BMW_FONT
            sub_para.font.size = Pt(16)
            sub_para.font.color.rgb = BMW_COLORS['black']
            sub_para.space_after = Pt(2)

    return slide


def create_diagram_slide(prs, title, diagram_items):
    """Create a slide with visual diagram boxes (up to 4)."""
    try:
        layout = _find_layout_by_name(prs, "Grid | 1")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Emu(488947), Emu(347184), Emu(11224684), Emu(400110)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(26)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    num_items = len(diagram_items[:4])
    colors = [BMW_COLORS['ocean_blue'], BMW_COLORS['accent_teal'],
              BMW_COLORS['blue_gray'], BMW_COLORS['pale_blue']]

    total_gap = 0.25
    total_width = 12.5
    box_width_val = (total_width - total_gap * (num_items - 1)) / num_items
    box_width = Inches(box_width_val)
    box_height = Inches(5.0)
    start_x = Inches(0.4)
    start_y = Inches(1.1)
    gap = Inches(total_gap)

    for i, (item_title, item_content) in enumerate(diagram_items[:4]):
        x_pos = start_x + (box_width + gap) * i

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, start_y,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Number circle
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x_pos + Inches(0.15), start_y + Inches(0.15),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = BMW_COLORS['white']
        num_circle.line.fill.background()

        num_box = slide.shapes.add_textbox(
            x_pos + Inches(0.15), start_y + Inches(0.2),
            Inches(0.5), Inches(0.4)
        )
        num_para = num_box.text_frame.paragraphs[0]
        num_para.text = str(i + 1)
        num_para.font.name = BMW_FONT
        num_para.font.size = Pt(18)
        num_para.font.bold = True
        num_para.font.color.rgb = colors[i % len(colors)]
        num_para.alignment = PP_ALIGN.CENTER

        # Item title
        item_title_box = slide.shapes.add_textbox(
            x_pos + Inches(0.15), start_y + Inches(0.75),
            box_width - Inches(0.3), Inches(0.7)
        )
        itf = item_title_box.text_frame
        itf.word_wrap = True
        ipara = itf.paragraphs[0]
        ipara.text = item_title.upper()
        ipara.font.name = BMW_FONT
        ipara.font.size = Pt(14)
        ipara.font.bold = True
        ipara.font.color.rgb = BMW_COLORS['white']

        # Item content
        content_box = slide.shapes.add_textbox(
            x_pos + Inches(0.15), start_y + Inches(1.5),
            box_width - Inches(0.3), box_height - Inches(1.7)
        )
        ctf = content_box.text_frame
        ctf.word_wrap = True
        cpara = ctf.paragraphs[0]
        cpara.text = item_content
        cpara.font.name = BMW_FONT
        cpara.font.size = Pt(12)
        cpara.font.color.rgb = BMW_COLORS['white']

    return slide


def create_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Create a two-column comparison slide."""
    try:
        layout = _find_layout_by_name(prs, "Grid | 1")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Emu(488947), Emu(347184), Emu(11224684), Emu(400110)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(26)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1),
        Inches(5.9), Inches(0.55)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = BMW_COLORS['ocean_blue']
    left_header.line.fill.background()

    lh_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.15), Inches(5.5), Inches(0.45)
    )
    lh_para = lh_box.text_frame.paragraphs[0]
    lh_para.text = left_title.upper()
    lh_para.font.name = BMW_FONT
    lh_para.font.size = Pt(16)
    lh_para.font.bold = True
    lh_para.font.color.rgb = BMW_COLORS['white']

    for i, point in enumerate(left_points[:5]):
        point_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.85 + i * 0.95), Inches(5.5), Inches(0.9)
        )
        ptf = point_box.text_frame
        ptf.word_wrap = True
        ppara = ptf.paragraphs[0]
        ppara.text = f"\u25aa  {point}"
        ppara.font.name = BMW_FONT
        ppara.font.size = Pt(14)
        ppara.font.color.rgb = BMW_COLORS['black']

    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.1),
        Inches(5.9), Inches(0.55)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = BMW_COLORS['blue_gray']
    right_header.line.fill.background()

    rh_box = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.15), Inches(5.5), Inches(0.45)
    )
    rh_para = rh_box.text_frame.paragraphs[0]
    rh_para.text = right_title.upper()
    rh_para.font.name = BMW_FONT
    rh_para.font.size = Pt(16)
    rh_para.font.bold = True
    rh_para.font.color.rgb = BMW_COLORS['white']

    for i, point in enumerate(right_points[:5]):
        point_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(1.85 + i * 0.95), Inches(5.5), Inches(0.9)
        )
        ptf = point_box.text_frame
        ptf.word_wrap = True
        ppara = ptf.paragraphs[0]
        ppara.text = f"\u25aa  {point}"
        ppara.font.name = BMW_FONT
        ppara.font.size = Pt(14)
        ppara.font.color.rgb = BMW_COLORS['black']

    return slide


def create_skill_slide(prs, skill_name, description, what_it_does, demo_output):
    """Create a skill showcase slide with description and demo output box."""
    try:
        layout = _find_layout_by_name(prs, "Grid | 1")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Title
    title_box = slide.shapes.add_textbox(
        Emu(488947), Emu(347184), Emu(11224684), Emu(400110)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = f"SKILL: {skill_name.upper()}"
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(26)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    # Skill badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.54), Inches(1.0),
        Inches(2.2), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = BMW_COLORS['ocean_blue']
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].text = "AGENT SKILL"
    badge_tf.paragraphs[0].font.name = BMW_FONT
    badge_tf.paragraphs[0].font.size = Pt(12)
    badge_tf.paragraphs[0].font.bold = True
    badge_tf.paragraphs[0].font.color.rgb = BMW_COLORS['white']
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(
        Inches(0.54), Inches(1.55), Inches(5.5), Inches(0.5)
    )
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    desc_para = desc_tf.paragraphs[0]
    desc_para.text = description
    desc_para.font.name = BMW_FONT
    desc_para.font.size = Pt(14)
    desc_para.font.italic = True
    desc_para.font.color.rgb = BMW_COLORS['accent_teal']

    # What it does — left column
    what_box = slide.shapes.add_textbox(
        Inches(0.54), Inches(2.15), Inches(5.5), Inches(4.5)
    )
    wtf = what_box.text_frame
    wtf.word_wrap = True

    for i, point in enumerate(what_it_does[:5]):
        if i == 0:
            para = wtf.paragraphs[0]
        else:
            para = wtf.add_paragraph()
        para.text = f"\u25aa  {point}"
        para.font.name = BMW_FONT
        para.font.size = Pt(14)
        para.font.color.rgb = BMW_COLORS['black']
        para.space_after = Pt(4)

    # Demo output box — right column
    demo_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.0),
        Inches(6.2), Inches(5.7)
    )
    demo_bg.fill.solid()
    demo_bg.fill.fore_color.rgb = BMW_COLORS['ice_blue_5']
    demo_bg.line.color.rgb = BMW_COLORS['blue_gray']
    demo_bg.line.width = Pt(1)

    # Demo label
    demo_label = slide.shapes.add_textbox(
        Inches(6.7), Inches(1.1), Inches(5.8), Inches(0.35)
    )
    dl_para = demo_label.text_frame.paragraphs[0]
    dl_para.text = "\u25b6  DEMO OUTPUT"
    dl_para.font.name = BMW_FONT
    dl_para.font.size = Pt(12)
    dl_para.font.bold = True
    dl_para.font.color.rgb = BMW_COLORS['ocean_blue']

    # Demo content
    demo_content = slide.shapes.add_textbox(
        Inches(6.7), Inches(1.5), Inches(5.8), Inches(5.0)
    )
    dtf = demo_content.text_frame
    dtf.word_wrap = True
    dpara = dtf.paragraphs[0]
    dpara.text = demo_output
    dpara.font.name = BMW_FONT
    dpara.font.size = Pt(10)
    dpara.font.color.rgb = BMW_COLORS['black']

    return slide


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    """Generate the AI for POs with GitHub Copilot presentation."""
    prs = Presentation(TEMPLATE_PATH)
    _delete_existing_slides(prs)
    _update_master_footer(prs, "DE-841 | March 2026 | AI for POs with GitHub Copilot")

    # ── Slide 1: Title ──────────────────────────────────────────────────
    create_title_slide(
        prs,
        "AI for Product Owners\nwith GitHub Copilot",
        "How I Use AI as a PO | Agent Skills | MCP Servers",
        "DE-841"
    )

    # ── Slide 2: Why AI for POs ─────────────────────────────────────────
    create_content_slide(
        prs,
        "Why Should POs Care About GitHub Copilot?",
        [
            "POs spend 60%+ of time on repetitive tasks: writing stories, tracking burndown, creating reports",
            "GitHub Copilot is not just for developers \u2014 it is an AI assistant for ANY knowledge work",
            "Custom Agent Skills teach Copilot YOUR domain: BMW processes, JIRA workflows, DLT analysis",
            "MCP Servers connect Copilot to BMW tools: JIRA, Confluence, GitHub, DLT logs",
            "Result: Automate hours of manual work into minutes with consistent, brand-compliant output"
        ]
    )

    # ── Slide 3: The agentSkills Repository ─────────────────────────────
    create_content_with_sub_bullets(
        prs,
        "My AI Toolkit: The agentSkills Repository",
        [
            ("Repository Structure \u2014 7 custom skills for daily PO work",
             ["Each skill has SKILL.md (instructions), references/, examples/",
              "Skills are reusable across Copilot, Claude, and any AI agent"]),
            ("How It Works \u2014 Skills + MCP = Complete Automation",
             ["SKILL.md teaches the AI what to do and how to do it",
              "MCP servers provide live access to JIRA, Confluence, GitHub, DLT logs"]),
            ("Key Principle \u2014 Build once, reuse everywhere",
             ["Write a skill once, use it in VS Code, CLI, or API-based agents",
              "Skills evolve over time with refinements and new reference data"]),
        ]
    )

    # ── Slide 4: Architecture Overview ──────────────────────────────────
    create_diagram_slide(
        prs,
        "Architecture: Skills + MCP Servers",
        [
            ("Agent Skills\n(SKILL.md)",
             "Domain instructions that\nteach AI your workflows.\n\n"
             "\u2022 bmw-brand-skill\n"
             "\u2022 burndown-tracker\n"
             "\u2022 dlt-defect-analyser\n"
             "\u2022 dlt-log-analyser\n"
             "\u2022 hlf-writer\n"
             "\u2022 user-story-writer"),
            ("MCP Servers\n(Tool Layer)",
             "Connect AI to live systems.\n\n"
             "\u2022 JIRA MCP \u2014 create/query tickets\n"
             "\u2022 Confluence MCP \u2014 read/write pages\n"
             "\u2022 GitHub MCP \u2014 PRs, repos, code\n"
             "\u2022 DLT MCP \u2014 parse vehicle logs\n\n"
             "Open standard: modelcontextprotocol.io"),
            ("GitHub Copilot\n(AI Engine)",
             "Orchestrates skills + tools.\n\n"
             "\u2022 Reads SKILL.md instructions\n"
             "\u2022 Calls MCP tools as needed\n"
             "\u2022 Generates code, reports, PPTs\n"
             "\u2022 Follows BMW brand guidelines\n\n"
             "Works in VS Code natively."),
            ("Output\n(Deliverables)",
             "Production-ready artifacts.\n\n"
             "\u2022 BMW-branded presentations\n"
             "\u2022 Burndown Excel reports\n"
             "\u2022 JIRA user stories & HLFs\n"
             "\u2022 DLT defect analysis reports\n\n"
             "Consistent quality every time."),
        ]
    )

    # ── Slide 5: Skill — BMW Brand Skill ────────────────────────────────
    create_skill_slide(
        prs,
        "BMW Brand Skill",
        "Ensures brand consistency by applying BMW styling patterns to all presentations and reports.",
        [
            "Extracts styling from PPT templates or uses BMW defaults",
            "Applies Ocean Blue #035970 theme, BMW typography, logos",
            "Generates python-pptx code with correct layouts and footers",
            "Validates brand compliance: colors, fonts, slide structure",
            "Supports both default (Presentation1.pptx) and custom templates"
        ],
        "Prompt: \"Create a presentation on AI for POs\"\n\n"
        "Output: AI_for_POs_BMW.pptx\n"
        "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
        "Slide 1: Title | Full Picture layout\n"
        "  \u2022 BMW gradient background (auto)\n"
        "  \u2022 Roundel + GROUP logos (auto)\n"
        "  \u2022 Department badge: DE-841\n\n"
        "Slide 2-11: Content slides\n"
        "  \u2022 26pt UPPERCASE Ocean Blue titles\n"
        "  \u2022 18pt body with \u25aa bullet style\n"
        "  \u2022 Master footer: DE-841 | Date | Topic\n"
        "  \u2022 Diagrams with brand color palette\n\n"
        "Brand validated \u2713 | 11 slides generated"
    )

    # ── Slide 6: Skill — Burndown Tracker ───────────────────────────────
    create_skill_slide(
        prs,
        "Burndown Tracker",
        "Tracks story point burndown vs forecasted velocity. Generates Excel with charts.",
        [
            "Extracts JIRA issues via JQL using Bearer token auth",
            "Groups issues by month, calculates cumulative burn",
            "Generates Excel: Monthly Burndown, Current Month, Closed Issues",
            "Creates burndown chart: actual (blue) vs forecast (red)",
            "Handles cross-year rollover (Dec tickets \u2192 Jan)"
        ],
        "Prompt: \"Generate burndown for label\n"
        "tsp_gl_robin2_eoy_2025, 1008 SP, 12 months\"\n\n"
        "Output: burndown_report.xlsx\n"
        "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
        "BURNDOWN SUMMARY \u2014 March 2026\n"
        "  Total forecasted:   1008 SP\n"
        "  Monthly forecast:   84 SP/month\n"
        "  Actual closed:      112 SP\n"
        "  Variance:           -140 BEHIND\n\n"
        "  Jan 2026: 29 SP  (forecast: 84)\n"
        "  Feb 2026: 83 SP  (forecast: 168)\n"
        "  Mar 2026:  0 SP  (forecast: 252)\n\n"
        "  Current Month: 3 In Review (22 SP),\n"
        "  6 In Specification, 4 New, 1 Open"
    )

    # ── Slide 7: Skill — DLT Defect Analyser ───────────────────────────
    create_skill_slide(
        prs,
        "DLT Defect Analyser",
        "Analyzes DLT logs from automotive ECUs to identify root causes of defects using MCP tools.",
        [
            "Uses DLT MCP server: stats, search, filter, context, time_range",
            "Pre-configured for Provisioning domain (PROV, PROI, PDS)",
            "Detects FATAL, ERROR, WARNING patterns with priority ranking",
            "Correlates events: gets context before/after each error",
            "Outputs prioritized root cause analysis with error codes"
        ],
        "Prompt: \"Provisioning stuck in DAS mode.\n"
        "Check /dltLogs/*.dlt\"\n\n"
        "Output: DLT Defect Analysis Report\n"
        "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
        "Top 3 Probable Causes:\n\n"
        "1. DAS File Unavailable [HIGH]\n"
        "   Code: 0x10d03000\n"
        "   \"DAS file is unavailable (index 12)\"\n"
        "   Action: Verify ECU coding\n\n"
        "2. DAS Index Not Obtained [HIGH]\n"
        "   Code: 0x10401000\n"
        "   Action: Re-code ECU with valid index\n\n"
        "3. No OTA Data Available [MEDIUM]\n"
        "   Expected until new OTA received"
    )

    # ── Slide 8: Skill — DLT Log Analyser ──────────────────────────────
    create_skill_slide(
        prs,
        "DLT Log Analyser",
        "Generic, domain-agnostic DLT analyzer. Interactively asks for your domain configuration.",
        [
            "Works for ANY automotive domain: Infotainment, ADAS, Body Control",
            "Asks for domain config: app IDs, ECU IDs, error patterns",
            "Uses same DLT MCP tools: stats, search, filter, context",
            "Supports local files, glob patterns, and GitHub repositories",
            "Provides domain-contextualized root cause analysis"
        ],
        "Prompt: \"Analyze /vehicle/logs/*.dlt\n"
        "Issue: Display black after update\"\n\n"
        "Skill asks: Domain? App IDs? ECU IDs?\n"
        "User: Infotainment, HMI/DISP/NAV, IHUD\n"
        "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
        "Analysis Report:\n"
        "  Domain: Infotainment\n"
        "  Found: 245 errors, 12 fatal\n\n"
        "1. Framebuffer Alloc Failed [CRITICAL]\n"
        "   DISP: Code 0xDEADFB01\n"
        "   Available: 128KB, Required: 8MB\n\n"
        "2. Update Service Restart Failed\n"
        "   HMI service not restarted properly"
    )

    # ── Slide 9: Skill — HLF Writer ────────────────────────────────────
    create_skill_slide(
        prs,
        "HLF Writer",
        "Co-author skill to help POs write High Level Features for JIRA with all required sections.",
        [
            "Guides through structured HLF sections step by step",
            "Sections: Problem Statement, Goals, Requirements, NFRs, Dependencies",
            "Asks probing questions for NFRs: error cases, robustness, corner cases",
            "Supports ICONPM, NODE0PM, IDCEVOPM project creation via JIRA MCP",
            "Follows BMW telematics domain conventions"
        ],
        "Prompt: \"Write an HLF for provisioning\n"
        "hot-reload feature\"\n\n"
        "Output: JIRA HLF ticket\n"
        "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
        "h3. Problem Statement\n"
        "  Current: Full restart for OTA config\n"
        "  Impact: 500K vehicles, 15min downtime\n\n"
        "h3. Goals and Objectives\n"
        "  Reduce update time: 15min \u2192 30sec\n"
        "  Zero service interruption\n\n"
        "h3. Non-Functional Requirements\n"
        "  \u2022 Atomic config swap on failure\n"
        "  \u2022 Max 3 retries with backoff\n"
        "  \u2022 ENG vs PROD mode handling"
    )

    # ── Slide 10: Skill — User Story Writer ─────────────────────────────
    create_skill_slide(
        prs,
        "User Story Writer",
        "Creates well-structured user stories with Motivation, Info, Questions, Scope, and Acceptance Criteria.",
        [
            "Standard 5-section format: Motivation, Info, Questions, Scope, AC",
            "As a / I want / so that format with specific personas",
            "Includes standard + story-specific acceptance criteria",
            "Creates JIRA tickets in ICONSD or NODE0DEV via MCP after PO confirmation",
            "Supports YAML templates for bulk story creation via JIRA Tree"
        ],
        "Prompt: \"Write a story for VIN retry\n"
        "mechanism in provisioningd\"\n\n"
        "Output: JIRA Story\n"
        "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
        "Motivation:\n"
        "  As a provisioning service, I want to\n"
        "  retry VIN requests when VIN service\n"
        "  is unavailable so that provisioning\n"
        "  does not fail due to transient issues.\n\n"
        "Acceptance Criteria:\n"
        "  \u2022 Retry with exponential backoff (5x)\n"
        "  \u2022 DTC raised after all retries fail\n"
        "  \u2022 Unit tests for retry logic\n"
        "  \u2022 Code merged, docs updated"
    )

    # ── Slide 11: MCP Servers — The Glue ────────────────────────────────
    create_two_column_slide(
        prs,
        "MCP Servers: Connecting AI to BMW Systems",
        "What Are MCP Servers?",
        [
            "Open standard: Model Context Protocol (modelcontextprotocol.io)",
            "Think of it as USB-C for AI \u2014 universal tool connector",
            "AI reads, writes, and orchestrates across live systems",
            "BMW SW Factory provides MCP servers for internal tools",
            "No custom API code needed \u2014 just configure and connect"
        ],
        "MCP Servers I Use Daily",
        [
            "JIRA MCP \u2014 create stories, query issues, manage labels",
            "Confluence MCP \u2014 read/write pages, search content",
            "GitHub MCP \u2014 PRs, file content, repo exploration",
            "DLT MCP \u2014 parse/filter/search vehicle log files",
            "JIRA Tree MCP \u2014 bulk-create hierarchical issues from YAML"
        ]
    )

    # ── Slide 12: Key Takeaways & Next Steps ────────────────────────────
    create_content_slide(
        prs,
        "Key Takeaways and Next Steps",
        [
            "GitHub Copilot + Custom Skills = a PO's AI-powered co-pilot for daily work",
            "Start small: pick ONE repetitive task, write a SKILL.md, and let AI handle it",
            "MCP servers are the bridge \u2014 they connect AI to JIRA, Confluence, GitHub, and DLT logs",
            "Skills are shareable: build once, share with your team, evolve over time",
            "Next step: Try writing your first skill for a task you do every week"
        ]
    )

    # Save
    output_path = os.path.join(OUTPUT_DIR, "AI_for_POs_with_Copilot.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
