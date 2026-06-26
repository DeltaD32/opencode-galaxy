#!/usr/bin/env python3
"""
BMW Brand Skill - AI for POs Presentation Generator (TEMPLATE FLOW)
Creates a presentation using the PSPO-AI-Essentials_Overview.pptx template
and content sourced from the Confluence page "What the heck AI".
Department: DE-841

Usage:
    python examples/create_ai_for_pos_presentation.py

Output:
    AI_for_POs_BMW.pptx in the skill root folder
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
EXAMPLES_DIR = os.path.dirname(os.path.abspath(__file__))
SKILL_ROOT = os.path.dirname(EXAMPLES_DIR)
TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.dirname(SKILL_ROOT)),
                             "PSPO-AI-Essentials_Overview.pptx")

# Output path
OUTPUT_DIR = SKILL_ROOT

# Template flow font (from brand-styling-guide.template.md)
BMW_FONT = "Franklin Gothic Medium"  # Fallback: "Arial"

# BMW Brand Colors (shared theme - BMW Group 21)
BMW_COLORS = {
    'ocean_blue': RGBColor(0x03, 0x59, 0x70),
    'night_blue': RGBColor(0x17, 0x3B, 0x68),
    'black': RGBColor(0x00, 0x00, 0x00),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x92, 0xA2, 0xBD),
    'accent_teal': RGBColor(0x54, 0x8D, 0x9E),
    'light_teal': RGBColor(0x85, 0xAC, 0xB9),
    'pale_blue': RGBColor(0xAB, 0xC4, 0xCF),
    'blue_gray': RGBColor(0x85, 0xAC, 0xB9),
    'ice_blue_4': RGBColor(0xC8, 0xD7, 0xE0),
    'ice_blue_5': RGBColor(0xDE, 0xE5, 0xEC),
    'ice_blue_6': RGBColor(0xE8, 0xEB, 0xF1),
    'cyan': RGBColor(0x07, 0x9E, 0xDA),
    'green': RGBColor(0x50, 0x81, 0x30),
    'orange': RGBColor(0xE9, 0x6D, 0x0C),
    'footer_gray': RGBColor(0x7A, 0x7A, 0x7A),
}


# ── Helper Functions ─────────────────────────────────────────────────────

def _update_master_footer(prs, footer_text):
    """Update the master slide's built-in FußzeileAU1 footer text."""
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.name == 'FußzeileAU1' or shape.name == 'Fu\u00dFzeileAU1':
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = ''
                shape.text_frame.paragraphs[0].runs[0].text = footer_text


def _find_layout_by_name(prs, name):
    """Find a slide layout by name across all slide masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise ValueError(f"Layout '{name}' not found in any slide master")


def _delete_existing_slides(prs):
    """Remove all existing slides, keeping masters/layouts intact."""
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId is None:
            for attr_name in sldId.attrib:
                if attr_name.endswith("}id") or attr_name == "r:id":
                    rId = sldId.attrib[attr_name]
                    break
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


# ── Slide Builders ───────────────────────────────────────────────────────

def create_title_slide(prs, title, subtitle, department="DE-841"):
    """Create a title slide using 'Title | Full Picture' layout."""
    layout = _find_layout_by_name(prs, "Title | Full Picture")
    slide = prs.slides.add_slide(layout)

    nsmap_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'

    for shape in slide.shapes:
        sp = shape._element
        ph_el = sp.find(f'.//{nsmap_p}ph')
        if ph_el is None:
            continue
        ph_type = ph_el.get('type')
        ph_idx = ph_el.get('idx')

        if ph_type == 'ctrTitle':
            shape.text_frame.paragraphs[0].text = title.upper()
            for p in shape.text_frame.paragraphs:
                p.font.name = BMW_FONT
                p.font.size = Pt(40)
                p.font.color.rgb = BMW_COLORS['white']
                p.alignment = PP_ALIGN.LEFT

        elif ph_type == 'subTitle' and ph_idx == '1':
            shape.text_frame.paragraphs[0].text = subtitle.upper()
            for p in shape.text_frame.paragraphs:
                p.font.name = BMW_FONT
                p.font.size = Pt(20)
                p.font.color.rgb = BMW_COLORS['white']
                p.alignment = PP_ALIGN.LEFT

        elif ph_type == 'body' and ph_idx == '22':
            shape.text_frame.paragraphs[0].text = department
            for p in shape.text_frame.paragraphs:
                p.font.name = BMW_FONT
                p.font.size = Pt(14)
                p.font.color.rgb = BMW_COLORS['white']
                p.alignment = PP_ALIGN.LEFT

        elif ph_type in ('pic', 'dgm'):
            pass  # Background image and SmartArt logos inherit from layout

    return slide


def create_content_slide(prs, title, bullet_points):
    """Create a content slide with bullet points (template flow layout)."""
    try:
        layout = _find_layout_by_name(prs, "1_Benutzerdefiniertes Layout")
    except ValueError:
        try:
            layout = _find_layout_by_name(prs, "Grid | 1")
        except ValueError:
            layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Title — template flow: 28pt
    title_box = slide.shapes.add_textbox(
        Inches(0.24), Inches(0.38), Inches(12.36), Inches(0.47)
    )
    title_box.text_frame.word_wrap = True
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(28)
    title_para.font.bold = False
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']
    title_para.alignment = PP_ALIGN.LEFT

    # Body — template flow: 20pt body, starting at y=0.92"
    content_box = slide.shapes.add_textbox(
        Inches(0.31), Inches(0.92), Inches(12.88), Inches(6.34)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points[:5]):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = f"▪  {point}"
        para.font.name = BMW_FONT
        para.font.size = Pt(20)
        para.font.color.rgb = BMW_COLORS['black']
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(6)

    return slide


def create_diagram_slide(prs, title, diagram_items):
    """Create a slide with visual diagram boxes (up to 4)."""
    try:
        layout = _find_layout_by_name(prs, "1_Benutzerdefiniertes Layout")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Inches(0.24), Inches(0.38), Inches(12.36), Inches(0.47)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    num_items = len(diagram_items[:4])
    colors = [BMW_COLORS['ocean_blue'], BMW_COLORS['accent_teal'],
              BMW_COLORS['blue_gray'], BMW_COLORS['pale_blue']]

    total_gap = 0.25
    total_width = 12.5
    box_width_val = (total_width - total_gap * (num_items - 1)) / num_items
    box_width = Inches(box_width_val)
    box_height = Inches(5.0)
    start_x = Inches(0.4)
    start_y = Inches(1.1)
    gap = Inches(total_gap)

    for i, (item_title, item_content) in enumerate(diagram_items[:4]):
        x_pos = start_x + (box_width + gap) * i

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, start_y,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Number circle
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x_pos + Inches(0.15), start_y + Inches(0.15),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = BMW_COLORS['white']
        num_circle.line.fill.background()

        num_box = slide.shapes.add_textbox(
            x_pos + Inches(0.15), start_y + Inches(0.2),
            Inches(0.5), Inches(0.4)
        )
        num_para = num_box.text_frame.paragraphs[0]
        num_para.text = str(i + 1)
        num_para.font.name = BMW_FONT
        num_para.font.size = Pt(18)
        num_para.font.bold = True
        num_para.font.color.rgb = colors[i % len(colors)]
        num_para.alignment = PP_ALIGN.CENTER

        # Item title
        item_title_box = slide.shapes.add_textbox(
            x_pos + Inches(0.15), start_y + Inches(0.75),
            box_width - Inches(0.3), Inches(0.7)
        )
        tf = item_title_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = item_title.upper()
        para.font.name = BMW_FONT
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = BMW_COLORS['white']

        # Item content
        content_box = slide.shapes.add_textbox(
            x_pos + Inches(0.15), start_y + Inches(1.5),
            box_width - Inches(0.3), box_height - Inches(1.7)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = item_content
        para.font.name = BMW_FONT
        para.font.size = Pt(12)
        para.font.color.rgb = BMW_COLORS['white']

    return slide


def create_three_box_slide(prs, title, items):
    """Create a slide with 3 boxes horizontally."""
    try:
        layout = _find_layout_by_name(prs, "1_Benutzerdefiniertes Layout")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Inches(0.24), Inches(0.38), Inches(12.36), Inches(0.47)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    colors = [BMW_COLORS['ocean_blue'], BMW_COLORS['accent_teal'],
              BMW_COLORS['blue_gray']]

    box_width = Inches(3.9)
    box_height = Inches(5.2)
    start_x = Inches(0.4)
    start_y = Inches(1.1)
    gap = Inches(0.3)

    for i, (item_title, item_content) in enumerate(items[:3]):
        x_pos = start_x + (box_width + gap) * i

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, start_y,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        item_title_box = slide.shapes.add_textbox(
            x_pos + Inches(0.2), start_y + Inches(0.3),
            box_width - Inches(0.4), Inches(0.7)
        )
        tf = item_title_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = item_title.upper()
        para.font.name = BMW_FONT
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = BMW_COLORS['white']

        content_box = slide.shapes.add_textbox(
            x_pos + Inches(0.2), start_y + Inches(1.1),
            box_width - Inches(0.4), box_height - Inches(1.4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = item_content
        para.font.name = BMW_FONT
        para.font.size = Pt(13)
        para.font.color.rgb = BMW_COLORS['white']

    return slide


def create_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Create a two-column comparison slide."""
    try:
        layout = _find_layout_by_name(prs, "1_Benutzerdefiniertes Layout")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Inches(0.24), Inches(0.38), Inches(12.36), Inches(0.47)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    # Left column
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1),
        Inches(5.9), Inches(0.55)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = BMW_COLORS['ocean_blue']
    left_header.line.fill.background()

    lh_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.15), Inches(5.5), Inches(0.45)
    )
    lh_para = lh_box.text_frame.paragraphs[0]
    lh_para.text = left_title.upper()
    lh_para.font.name = BMW_FONT
    lh_para.font.size = Pt(16)
    lh_para.font.bold = True
    lh_para.font.color.rgb = BMW_COLORS['white']

    for i, point in enumerate(left_points[:5]):
        point_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.85 + i * 0.95), Inches(5.5), Inches(0.9)
        )
        tf = point_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = f"▪  {point}"
        para.font.name = BMW_FONT
        para.font.size = Pt(14)
        para.font.color.rgb = BMW_COLORS['black']

    # Right column
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.1),
        Inches(5.9), Inches(0.55)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = BMW_COLORS['blue_gray']
    right_header.line.fill.background()

    rh_box = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.15), Inches(5.5), Inches(0.45)
    )
    rh_para = rh_box.text_frame.paragraphs[0]
    rh_para.text = right_title.upper()
    rh_para.font.name = BMW_FONT
    rh_para.font.size = Pt(16)
    rh_para.font.bold = True
    rh_para.font.color.rgb = BMW_COLORS['white']

    for i, point in enumerate(right_points[:5]):
        point_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(1.85 + i * 0.95), Inches(5.5), Inches(0.9)
        )
        tf = point_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = f"▪  {point}"
        para.font.name = BMW_FONT
        para.font.size = Pt(14)
        para.font.color.rgb = BMW_COLORS['black']

    return slide


def create_content_with_sub_bullets(prs, title, items):
    """Create a content slide with main bullets and sub-bullets.

    items: list of tuples (main_text, [sub_text, ...])
    """
    try:
        layout = _find_layout_by_name(prs, "1_Benutzerdefiniertes Layout")
    except ValueError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    title_box = slide.shapes.add_textbox(
        Inches(0.24), Inches(0.38), Inches(12.36), Inches(0.47)
    )
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title.upper()
    title_para.font.name = BMW_FONT
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = BMW_COLORS['ocean_blue']

    content_box = slide.shapes.add_textbox(
        Inches(0.31), Inches(0.92), Inches(12.88), Inches(6.34)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    first = True
    for main_text, sub_items in items:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.text = f"▪  {main_text}"
        para.font.name = BMW_FONT
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = BMW_COLORS['black']
        para.space_before = Pt(8)
        para.space_after = Pt(2)

        for sub in sub_items[:3]:
            sub_para = tf.add_paragraph()
            sub_para.text = f"     -  {sub}"
            sub_para.font.name = BMW_FONT
            sub_para.font.size = Pt(16)
            sub_para.font.color.rgb = BMW_COLORS['black']
            sub_para.space_after = Pt(2)

    return slide


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    """Generate the AI for POs presentation using the PSPO template."""
    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        _delete_existing_slides(prs)
    else:
        raise FileNotFoundError(
            f"Template not found: {TEMPLATE_PATH}\n"
            "Please provide PSPO-AI-Essentials_Overview.pptx"
        )

    _update_master_footer(prs, "DE-841 | February 2026 | AI for POs")

    # ── Slide 1: Title ──────────────────────────────────────────────────
    create_title_slide(
        prs,
        "AI for Product Owners",
        "What the Heck is AI? | A Practical Guide for POs, Developers & Testers",
        "DE-841"
    )

    # ── Slide 2: Executive Summary ──────────────────────────────────────
    create_content_slide(
        prs,
        "AI Can Be Broken Down into Three Practical Categories",
        [
            "Applied AI (Basic): Use existing AI tools like GAIA, GitHub Copilot, and BMW MCP servers in daily work",
            "Build with Off-the-Shelf Models (Intermediate): Create AI-powered workflows and agents without training models",
            "Build with Custom Models (Advanced): Train and fine-tune models for domain-specific use cases",
            "Key principle: Start simple — optimise single LLM calls before building agentic systems",
            "BMW provides internal tools (GAIA, Paper Tool, SnAPP) alongside open-source options"
        ]
    )

    # ── Slide 3: Three Categories Overview ──────────────────────────────
    create_three_box_slide(
        prs,
        "Three Categories of AI for POs",
        [
            ("Basic\nApplied AI",
             "I want to use available AI tools in my day-to-day work so that I am AI literate.\n\n"
             "• GitHub Copilot\n"
             "• GAIA (+ Browser Plugin)\n"
             "• Paper Tool\n"
             "• SnAPP (No-code)\n"
             "• Microsoft Copilot\n"
             "• BMW MCP Servers\n\n"
             "Goal: Become AI-literate, boost daily productivity."),
            ("Intermediate\nBuild with Existing Models",
             "I want to build applications that use AI to automate my repetitive tasks.\n\n"
             "• Ollama + OpenWebUI (local LLMs)\n"
             "• n8n (no-code workflows)\n"
             "• LangChain / LangGraph\n"
             "• CrewAI (multi-agent)\n"
             "• BMW GAIA LLM APIs\n\n"
             "Goal: Automate workflows without training models."),
            ("Advanced\nBuild with Custom Models",
             "I want to build applications that need models trained on my own data.\n\n"
             "• GAIA custom apps\n"
             "• CAIP (Connected AI Platform)\n"
             "• Hugging Face fine-tuning\n"
             "• Unsloth\n"
             "• RAG architecture\n\n"
             "Goal: Domain-specific AI with proprietary data.")
        ]
    )

    # ── Slide 4: Agents & Workflows Explained ───────────────────────────
    create_two_column_slide(
        prs,
        "Before We Start: Workflows vs Agents",
        "Workflows",
        [
            "LLMs and tools orchestrated through predefined code paths",
            "Predictable and consistent for well-defined tasks",
            "Lower latency and cost than full agents",
            "Best for structured, repeatable processes",
            "Example: Automated PR review with fixed steps"
        ],
        "Agents",
        [
            "LLMs dynamically direct their own processes and tool usage",
            "Flexible, model-driven decision-making at scale",
            "Higher latency and cost — use only when needed",
            "Best when flexibility outweighs predictability",
            "Tip: Don't make everything an agent — start simple!"
        ]
    )

    # ── Slide 5: Deep Dive — Applied AI (BMW Tools) ─────────────────────
    create_content_with_sub_bullets(
        prs,
        "Applied AI: BMW Internal Tools",
        [
            ("GitHub Copilot — Coding Agent",
             ["AI pair-programmer integrated into VS Code",
              "Code suggestions, explanations, and refactoring"]),
            ("GAIA — BMW's Generative AI Assistant",
             ["Chat interface + Browser Plugin for quick answers",
              "Pinned apps: Provisioning User Support and more"]),
            ("Paper Tool (paper.bmwgroup.net)",
             ["Chat with BMW data: Features, Requirements, Defects",
              "Works best with exact ticket numbers for summarisation"]),
            ("SnAPP — No-Code AI Application Builder",
             ["Build AI workflows without writing code",
              "Use Case Store with pre-built templates"]),
            ("BMW MCP Servers — Tool Integration Layer",
             ["MCP servers for Jira, Confluence, GitHub, Codebeamer",
              "Provided by SW Factory: enables AI tool orchestration"]),
        ]
    )

    # ── Slide 6: Deep Dive — Applied AI (Examples) ──────────────────────
    create_content_slide(
        prs,
        "Applied AI in Action: Practical Examples",
        [
            "Automate workflows for devs/testers/POs — creating PPTs, writing user stories, writing HLFs, pre-DLT analysis",
            "Vibe-code to build new development apps and prototypes — dashboards for test results, quick PoCs",
            "Use GAIA Browser Plugin for in-context answers without leaving your browser",
            "Leverage MCP servers to connect AI to Jira, Confluence, GitHub in a single workflow",
            "BMW AI Learning Path available for structured upskilling (SharePoint)"
        ]
    )

    # ── Slide 7: Deep Dive — Build with Off-the-Shelf Models ────────────
    create_diagram_slide(
        prs,
        "Intermediate: Build AI Apps Without Training Models",
        [
            ("Run LLMs Locally",
             "Don't send BMW data to the cloud.\n\n"
             "• Ollama: Run open-source models locally\n"
             "• OpenWebUI: Chat interface for Ollama\n"
             "• BMW GAIA LLM APIs (cloud but BMW-hosted)\n\n"
             "Privacy-first AI for PoCs."),
            ("No-Code Workflows",
             "Build AI workflows visually.\n\n"
             "• n8n: Self-hosted workflow automation\n"
             "• Pre-built templates available\n"
             "• Fair-code licence\n\n"
             "Ideal for POs and testers."),
            ("Coding Frameworks",
             "For developers building agents.\n\n"
             "• LangChain: Popular agent framework\n"
             "• LangGraph: Step-by-step orchestration\n"
             "• CrewAI: Multi-agent collaboration\n\n"
             "Evaluate needs first, then pick."),
            ("Key Principle",
             "Start with the simplest solution.\n\n"
             "• Optimise single LLM calls first\n"
             "• Add retrieval + in-context examples\n"
             "• Only then build agentic systems\n\n"
             "Complexity ≠ better results.")
        ]
    )

    # ── Slide 8: Deep Dive — Custom Models & RAG ────────────────────────
    create_content_slide(
        prs,
        "Advanced: Custom Models and RAG Architecture",
        [
            "Fine-tuning: Use Hugging Face Transformers or Unsloth to train models on BMW-specific data",
            "GAIA Custom Apps: No-code way to build apps with your own data inside BMW's ecosystem",
            "CAIP (Connected AI Platform): BMW's vision for managed custom model training and deployment",
            "RAG (Retrieval Augmented Generation): Enhance LLMs with external BMW knowledge without full re-training",
            "RAG = cost-effective alternative to fine-tuning — retrieve relevant docs, then generate grounded answers"
        ]
    )

    # ── Slide 9: RAG Architecture Deep Dive ─────────────────────────────
    create_diagram_slide(
        prs,
        "RAG Architecture: Key Components",
        [
            ("Vector Databases",
             "Store and search document embeddings.\n\n"
             "• Chroma — quick prototyping\n"
             "• Weaviate — production, hybrid search\n"
             "• Qdrant — high-throughput\n"
             "• Milvus — enterprise scale\n"
             "• FAISS — research, similarity search"),
            ("Embedding Models",
             "Convert text to vector representations.\n\n"
             "• Sentence Transformers\n"
             "• OpenAI Embeddings\n"
             "• BGE (BAAI)\n"
             "• Instructor Embeddings\n\n"
             "Choose based on language & domain."),
            ("RAG Frameworks",
             "Orchestrate retrieval + generation.\n\n"
             "• LangChain — most popular\n"
             "• LlamaIndex — advanced indexing\n"
             "• Haystack — production-ready\n\n"
             "All support Python; LangChain\nalso supports TypeScript."),
            ("Doc Processing",
             "Parse documents for the pipeline.\n\n"
             "• Unstructured — multi-format\n"
             "• LangChain Loaders — 100+ types\n"
             "• PyPDF2 / PDFMiner\n"
             "• python-docx\n\n"
             "Feed clean data into your RAG.")
        ]
    )

    # ── Slide 10: MCP Protocol & Agent Skills ───────────────────────────
    create_two_column_slide(
        prs,
        "Enabling Technologies: MCP & Agent Skills",
        "MCP Protocol",
        [
            "Open standard for connecting AI to external tools",
            "BMW MCP servers: Jira, Confluence, GitHub, Codebeamer",
            "Enables AI to read, write, and orchestrate across systems",
            "Think of it as a USB-C for AI — universal tool connector",
            "Docs: modelcontextprotocol.io"
        ],
        "Agent Skills",
        [
            "Reusable instruction sets that teach AI specific tasks",
            "Build once, reuse across Claude, Copilot, and APIs",
            "Skills + MCP = complete AI automation solution",
            "Examples: PPT creation, user story writing, DLT analysis",
            "Docs: agentskills.io"
        ]
    )

    # ── Slide 11: Recommended Next Steps ────────────────────────────────
    create_content_slide(
        prs,
        "Recommended Next Steps for Your Team",
        [
            "Week 1: Get AI-literate — set up GAIA Browser Plugin, explore Paper Tool and GitHub Copilot",
            "Week 2-3: Identify 2-3 repetitive tasks and prototype automation with n8n or Agent Skills",
            "Month 2: Evaluate if your use case needs custom models or RAG — start with Ollama + Chroma locally",
            "Month 3: Present findings to leadership with a working PoC and clear cost-benefit analysis",
            "Ongoing: Follow BMW AI Learning Path and contribute skills back to the team"
        ]
    )

    # Save
    output_path = os.path.join(OUTPUT_DIR, "AI_for_POs_BMW.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
