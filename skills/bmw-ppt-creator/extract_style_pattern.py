#!/usr/bin/env python3
"""
Extract basic style patterns from a PowerPoint file.
Outputs a JSON summary of slide dimensions, title/body font info, and detected colors.

Usage:
  python extract_style_pattern.py /path/to/presentation.pptx [output.json]
"""

import json
import sys
from collections import Counter
from pptx import Presentation


def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def get_color(font):
    if font is None or font.color is None:
        return None
    try:
        if font.color.rgb is None:
            return None
        return font.color.rgb
    except Exception:
        return None


def extract_text_styles_from_shapes(shapes, source):
    styles = []
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            if para.runs:
                for run in para.runs:
                    font = run.font or para.font
                    color = get_color(font)
                    styles.append({
                        "text": run.text.strip(),
                        "font": font.name,
                        "size": float(font.size.pt) if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic,
                        "color": rgb_to_hex(color),
                        "source": source,
                    })
            else:
                font = para.font
                color = get_color(font)
                styles.append({
                    "text": para.text.strip() if para.text else "",
                    "font": font.name,
                    "size": float(font.size.pt) if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                    "color": rgb_to_hex(color),
                    "source": source,
                })
    return styles


def summarize_font_sizes(styles):
    sizes = [s["size"] for s in styles if s["size"]]
    return dict(Counter(sizes))


def summarize_fonts(styles):
    fonts = [s["font"] for s in styles if s["font"]]
    return dict(Counter(fonts))


def summarize_colors(styles):
    colors = [s["color"] for s in styles if s["color"]]
    return dict(Counter(colors))


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_style_pattern.py /path/to/presentation.pptx [output.json]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    prs = Presentation(pptx_path)

    all_styles = []
    for slide in prs.slides:
        all_styles.extend(extract_text_styles_from_shapes(slide.shapes, "slide"))

    master_styles = []
    layout_styles = []
    for master in prs.slide_masters:
        master_styles.extend(extract_text_styles_from_shapes(master.shapes, "master"))
        for layout in master.slide_layouts:
            layout_styles.extend(extract_text_styles_from_shapes(layout.shapes, "layout"))

    report = {
        "slide_dimensions": {
            "width_in": prs.slide_width / 914400,
            "height_in": prs.slide_height / 914400,
        },
        "fonts": summarize_fonts(all_styles + master_styles + layout_styles),
        "font_sizes_pt": summarize_font_sizes(all_styles + master_styles + layout_styles),
        "colors": summarize_colors(all_styles + master_styles + layout_styles),
        "sample_text_styles": [s for s in (master_styles + layout_styles + all_styles) if s["text"]][:30],
        "master_text_styles_count": len(master_styles),
        "layout_text_styles_count": len(layout_styles),
        "slide_text_styles_count": len(all_styles),
    }

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(report, f, indent=2)
        print(f"Saved style report to {output_path}")
    else:
        print(json.dumps(report, indent=2))


if __name__ == "__main__":
    main()
